from io import BytesIO
import math
import pandas as pd
import datetime
import subprocess
import os
from pptx import Presentation
from pptx.util import Inches, Pt


COLUMNAS_REPORTE = [
    "RUT Cliente",
    "Cliente",
    "Ejecutivo",
    "ID Deudor",
    "Deudor",
    "Fecha Otorgamiento",
    "Tipo Documento",
    "N°Documento",
    "Fecha Vencimiento",
    "Días Mora",
    "Monto Documento",
    "Monto Recaudado",
    "Capital Amortizado",
    "Monto Saldo",
]

FILAS_POR_SLIDE = 8
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

TABLE_WIDTH = Inches(12.3)
TABLE_HEIGHT = Inches(5.2)

def pptx_a_pdf(pptx_bytes: bytes, output_pdf_path: str) -> None:
    """
    Convierte un PPTX en memoria a PDF usando LibreOffice headless.
    Guarda el PDF en output_pdf_path.
    """
    tmp_pptx = "temp_reporte.pptx"
    with open(tmp_pptx, "wb") as f:
        f.write(pptx_bytes)

    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", os.path.dirname(output_pdf_path),
        tmp_pptx
    ], check=True)

    os.rename(os.path.splitext(tmp_pptx)[0] + ".pdf", output_pdf_path)

    os.remove(tmp_pptx)

def generar_ppt(excel_bytes: bytes) -> bytes:
    df = pd.read_excel(BytesIO(excel_bytes), header=1)

    columnas_faltantes = [c for c in COLUMNAS_REPORTE if c not in df.columns]
    if columnas_faltantes:
        raise ValueError(f"Faltan columnas: {', '.join(columnas_faltantes)}")

    df = df[COLUMNAS_REPORTE]

    prs = Presentation("app/template.pptx")

    cliente = str(df.iloc[0]["Cliente"])
    rut_cliente = str(df.iloc[0]["RUT Cliente"])

    portada = prs.slides[0]
    for shape in portada.shapes:
        if shape.name == "info_cliente" and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()

            p1 = tf.paragraphs[0]
            p1.text = f"{cliente}\n{rut_cliente}"
            p1.font.bold = True
            p1.font.size = Pt(20)

    COLUMNAS_TABLA = COLUMNAS_REPORTE[3:]
    total_slides = math.ceil(len(df) / FILAS_POR_SLIDE)

    columnas_totales = ["Monto Documento", "Monto Recaudado", "Capital Amortizado", "Monto Saldo"]
    totales_globales = {col: df[col].sum() for col in columnas_totales}

    for slide_idx in range(total_slides):
        inicio = slide_idx * FILAS_POR_SLIDE
        fin = inicio + FILAS_POR_SLIDE
        df_slice = df.iloc[inicio:fin]

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        es_ultima_slide = (slide_idx == total_slides - 1)
        filas = len(df_slice) + 1 + (1 if es_ultima_slide else 0)
        columnas = len(COLUMNAS_TABLA)

        left = (SLIDE_WIDTH - TABLE_WIDTH) // 2
        top = (SLIDE_HEIGHT - TABLE_HEIGHT) // 4

        table = slide.shapes.add_table(
            filas,
            columnas,
            left,
            top,
            TABLE_WIDTH,
            TABLE_HEIGHT,
        ).table

        ancho_columna = int(TABLE_WIDTH / columnas)
        alto_fila = int(TABLE_HEIGHT / FILAS_POR_SLIDE)
        for col in table.columns:
            col.width = ancho_columna
        for row in table.rows:
            row.height = alto_fila

        for col_idx, col_name in enumerate(COLUMNAS_TABLA):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(8)

        for row_idx, row in enumerate(df_slice[COLUMNAS_TABLA].itertuples(index=False), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)

                if pd.isna(value):
                    texto = ""
                elif COLUMNAS_TABLA[col_idx] in ["Cliente", "Deudor"]:
                    texto = str(value).title()
                elif COLUMNAS_TABLA[col_idx] in ["Fecha Otorgamiento", "Fecha Vencimiento"]:
                    if isinstance(value, (pd.Timestamp, datetime.datetime)):
                        texto = value.strftime("%d-%m-%Y")
                    else:
                        texto = str(value)
                else:
                    texto = str(value)

                cell.text = texto
                cell.text_frame.paragraphs[0].font.size = Pt(8)

        if es_ultima_slide:
            total_row_idx = len(df_slice) + 1
            for col_idx, col_name in enumerate(COLUMNAS_TABLA):
                cell = table.cell(total_row_idx, col_idx)
                if col_name in columnas_totales:
                    suma = totales_globales[col_name]
                    cell.text = f"{suma:,.2f}"
                elif col_idx == 0:
                    cell.text = "TOTAL"
                else:
                    cell.text = ""
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(8)

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return output.read()


